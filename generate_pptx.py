from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────────────
VIOLET    = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_D  = RGBColor(0x4C, 0x1D, 0x95)
BLANC     = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_C    = RGBColor(0xE5, 0xE7, 0xEB)
GRIS_F    = RGBColor(0x1F, 0x2A, 0x40)
NOIR      = RGBColor(0x0D, 0x0D, 0x14)
VERT      = RGBColor(0x10, 0xB9, 0x81)
OR        = RGBColor(0xF5, 0x9E, 0x0B)
ROUGE     = RGBColor(0xEF, 0x44, 0x44)

BLANK = prs.slide_layouts[6]  # blank layout

def bg(slide, couleur=NOIR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = couleur

def box(slide, left, top, width, height, text, fontsize=18, bold=False,
        color=BLANC, bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        fill = tf.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return tf

def rect(slide, left, top, width, height, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def bullet_box(slide, left, top, width, height, items, fontsize=15,
               color=BLANC, bg_color=None, spacing=1.2):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(fontsize)
        run.font.color.rgb = color
    return tf

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.4, VIOLET_D)
    rect(slide, 0, 1.35, 13.33, 0.07, VIOLET)
    box(slide, 0.4, 0.12, 12, 0.7, title, fontsize=30, bold=True,
        color=BLANC, align=PP_ALIGN.LEFT)
    if subtitle:
        box(slide, 0.4, 0.82, 12, 0.45, subtitle, fontsize=16,
            color=GRIS_C, align=PP_ALIGN.LEFT, italic=True)

def footer(slide, num, total=13):
    rect(slide, 0, 7.15, 13.33, 0.35, VIOLET_D)
    box(slide, 0.3, 7.17, 6, 0.28, "BorneInteract — BTS CIEL-IR 2025-2026 — Adictiz",
        fontsize=10, color=GRIS_C)
    box(slide, 11.5, 7.17, 1.5, 0.28, f"{num} / {total}",
        fontsize=10, color=GRIS_C, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
rect(s, 0, 0, 13.33, 7.5, NOIR)
rect(s, 0, 0, 13.33, 3.8, VIOLET_D)
rect(s, 0, 3.75, 13.33, 0.12, VIOLET)

box(s, 0.6, 0.4, 12, 0.9, "BorneInteract", fontsize=52, bold=True,
    color=BLANC, align=PP_ALIGN.LEFT)
box(s, 0.6, 1.35, 12, 0.65, "Borne interactive phygitale en magasin", fontsize=24,
    color=GRIS_C, italic=True)
box(s, 0.6, 2.1, 12, 0.55, "Projet BTS CIEL-IR · Partenaire : Adictiz · Session 2025-2026",
    fontsize=16, color=GRIS_C)

rect(s, 0.6, 4.1, 3.5, 1.3, GRIS_F)
box(s, 0.8, 4.2, 3.2, 0.4, "Étudiant 1", fontsize=14, color=GRIS_C, bold=True)
bullet_box(s, 0.8, 4.6, 3.2, 0.8, [
    "• Sécurité & BDD",
    "• Hardware IoT (RFID, TCS3200)",
    "• Jeu Memory · 2ème Chance · ESP32-CAM"
], fontsize=11, color=GRIS_C)

rect(s, 4.5, 4.1, 3.8, 1.3, GRIS_F)
box(s, 4.7, 4.2, 3.5, 0.4, "Équipe", fontsize=14, color=GRIS_C, bold=True)
bullet_box(s, 4.7, 4.6, 3.5, 0.8, [
    "DJELASSI Souhir",
    "BENGORA Bouguerra",
    "NEGGAOUI Ali"
], fontsize=11, color=GRIS_C)

rect(s, 8.7, 4.1, 3.8, 1.3, GRIS_F)
box(s, 8.9, 4.2, 3.5, 0.4, "Oral", fontsize=14, color=GRIS_C, bold=True)
bullet_box(s, 8.9, 4.6, 3.5, 0.8, [
    "Durée : 15 minutes",
    "6 missions présentées",
    "Démos possibles"
], fontsize=11, color=GRIS_C)

footer(s, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Contexte & Architecture
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Contexte & Architecture globale",
           "Une borne phygitale : ticket de caisse → jeu → gain")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Le parcours client", fontsize=14, bold=True, color=VIOLET)
bullet_box(s, 0.6, 2.15, 5.4, 4.3, [
    "1.  Client passe en caisse",
    "2.  Reçoit ticket avec QR code / code-barres",
    "3.  Scanne son ticket sur la borne",
    "4.  Joue un mini-jeu (Memory, Roue, Jackpot…)",
    "5.  Gagne une réduction ou un bon cadeau",
    "6.  En cas de perte → 2ème chance via smartphone",
    "",
    "Objectif Adictiz : générer des leads qualifiés",
    "en magasin & fidéliser les clients."
], fontsize=13, color=BLANC)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Stack technique", fontsize=14, bold=True, color=VIOLET)
bullet_box(s, 6.8, 2.15, 5.9, 4.3, [
    "Hardware :  ESP8266 + RC522 + TCS3200 + ESP32-CAM",
    "Backend  :  PHP 8.2 + MySQL 8.0 (Docker Compose)",
    "Frontend :  HTML5 / JS (borne) + Flutter (mobile)",
    "Réseau   :  WiFi 'labo_snir' → API REST",
    "Sécurité :  bcrypt · PDO prepared · CSRF · Rate limit",
    "",
    "3 étudiants, 3 périmètres complémentaires :",
    "  Étudiant 1 → Hardware + Sécurité + Jeu Memory",
    "  Étudiant 2 → App Flutter + Jeux web",
    "  Étudiant 3 → API PHP + BDD + Admin dashboard"
], fontsize=13, color=BLANC)

footer(s, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Mes 6 missions
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mes 6 missions", "Vue d'ensemble des réalisations")

missions = [
    ("🔒", "Sécurité BDD",    "bcrypt · PDO · CSRF · Rate limit"),
    ("📡", "RFID RC522",      "Badge client → API en WiFi"),
    ("🎨", "TCS3200 & Jeton", "Détection couleur jeton rouge"),
    ("🎮", "Jeu Memory",      "HTML5/JS · Fisher-Yates · Score"),
    ("🔄", "2ème Chance",     "QR code → rejouer gratuitement"),
    ("📷", "ESP32-CAM",       "Streaming MJPEG surveillance"),
]

cols = [(0.3, 1.6), (4.65, 1.6), (8.95, 1.6),
        (0.3, 4.15), (4.65, 4.15), (8.95, 4.15)]

for (icon, title, desc), (l, t) in zip(missions, cols):
    rect(s, l, t, 3.9, 2.3, GRIS_F)
    rect(s, l, t, 3.9, 0.55, VIOLET_D)
    box(s, l+0.15, t+0.08, 3.6, 0.42, f"{icon}  {title}",
        fontsize=16, bold=True, color=BLANC)
    box(s, l+0.15, t+0.65, 3.6, 1.5, desc,
        fontsize=13, color=GRIS_C)

footer(s, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Sécurité BDD
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 1 — Sécurité Base de Données",
           "Conformité ANSSI : zéro mot de passe en clair, zéro injection SQL")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Mesures de sécurité", fontsize=14, bold=True, color=VERT)
bullet_box(s, 0.6, 2.15, 5.4, 4.3, [
    "✅  Hachage bcrypt (coût adaptatif)",
    "    Jamais de mot de passe stocké en clair",
    "",
    "✅  Prepared Statements PDO",
    "    Aucune concaténation de chaîne dans les requêtes",
    "    Paramètres liés → SQL injection impossible",
    "",
    "✅  Tokens CSRF",
    "    Protection contre les requêtes forgées",
    "",
    "✅  Rate Limiting",
    "    Max 3 tentatives d'inscription par heure/IP",
    "",
    "✅  Logs complets",
    "    Timestamp + UID + résultat pour chaque scan"
], fontsize=12, color=BLANC)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Exemple — Prepared Statement PDO", fontsize=14, bold=True, color=VERT)
code = (
    "// ✗ DANGEREUX (injection SQL possible)\n"
    "$sql = \"SELECT * FROM clients\n"
    "        WHERE email='\" . $email . \"'\";\n\n"
    "// ✓ SÉCURISÉ (PDO prepared)\n"
    "$stmt = $pdo->prepare(\n"
    "    'SELECT * FROM clients\n"
    "     WHERE email = :email'\n"
    ");\n"
    "$stmt->execute([':email' => $email]);\n\n"
    "// ✓ Mot de passe (bcrypt)\n"
    "$hash = password_hash($mdp, PASSWORD_BCRYPT);\n"
    "password_verify($mdp, $hash_bdd);"
)
box(s, 6.8, 2.2, 5.9, 4.4, code, fontsize=11, color=VERT, bg_color=RGBColor(0x0A, 0x1A, 0x0A))

footer(s, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — RFID RC522
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 2 — Intégration RFID RC522",
           "Badge client → UID → API PHP → ouverture session jeu")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Câblage ESP8266 → RC522", fontsize=14, bold=True, color=OR)
cable = (
    "ESP8266          RC522\n"
    "─────────────────────────\n"
    "D2  (GPIO4)  →   SDA\n"
    "D5  (GPIO14) →   SCK\n"
    "D7  (GPIO13) →   MOSI\n"
    "D6  (GPIO12) →   MISO\n"
    "D1  (GPIO5)  →   RST\n"
    "3.3V         →   VCC\n"
    "GND          →   GND\n"
)
box(s, 0.6, 2.2, 5.4, 2.8, cable, fontsize=12, color=OR,
    bg_color=RGBColor(0x1A, 0x14, 0x00))

box(s, 0.6, 5.1, 5.4, 0.4, "Fonctionnement", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 0.6, 5.5, 5.4, 1.2, [
    "• Anti-rebond 3 s (évite double scan)",
    "• POST /api/rfid_scan.php  {\"uid\":\"XXXXX\"}",
    "• LED verte = accès OK  /  LED rouge = refus"
], fontsize=12, color=BLANC)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Code firmware (extrait)", fontsize=14, bold=True, color=OR)
code2 = (
    "String lireUID() {\n"
    "    String uid = \"\";\n"
    "    for (byte i = 0; i < rfid.uid.size; i++) {\n"
    "        if (rfid.uid.uidByte[i] < 0x10) uid += '0';\n"
    "        uid += String(rfid.uid.uidByte[i], HEX);\n"
    "    }\n"
    "    uid.toUpperCase();\n"
    "    return uid;\n"
    "}\n\n"
    "bool envoyerScan(const String& uid) {\n"
    "    HTTPClient http;\n"
    "    String url = \"http://\" + API_HOST\n"
    "               + \":\" + API_PORT\n"
    "               + \"/api/rfid_scan.php\";\n"
    "    http.begin(wifiClient, url);\n"
    "    String body = '{\"uid\":\"' + uid + '\"}';\n"
    "    int code = http.POST(body);\n"
    "    return (code == 200);\n"
    "}"
)
box(s, 6.8, 2.2, 5.9, 4.4, code2, fontsize=11, color=OR,
    bg_color=RGBColor(0x1A, 0x14, 0x00))

footer(s, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — TCS3200 & Jeton
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 3 — Capteur TCS3200 & Jeton physique",
           "Détecter un jeton rouge = interaction physique unique")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Câblage ESP8266 → TCS3200", fontsize=14, bold=True, color=ROUGE)
cable3 = (
    "ESP8266          TCS3200\n"
    "─────────────────────────\n"
    "D3  (GPIO0)  →   S2\n"
    "D4  (GPIO2)  →   S3\n"
    "D0  (GPIO16) →   OUT\n"
    "3.3V (fixé)  →   S0\n"
    "GND  (fixé)  →   S1\n"
    "GND          →   OE (actif bas)\n"
)
box(s, 0.6, 2.2, 5.4, 2.4, cable3, fontsize=12, color=ROUGE,
    bg_color=RGBColor(0x1A, 0x04, 0x04))

box(s, 0.6, 4.7, 5.4, 0.35, "Plages de calibrage (jeton rouge)", fontsize=12,
    bold=True, color=BLANC)
calib = (
    "Couleur   MIN    MAX\n"
    "──────────────────────\n"
    "Rouge      0     200   ← faible\n"
    "Vert     200     999   ← élevé\n"
    "Bleu     200     999   ← élevé\n"
)
box(s, 0.6, 5.1, 5.4, 1.5, calib, fontsize=12, color=ROUGE,
    bg_color=RGBColor(0x1A, 0x04, 0x04))

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Principe physique", fontsize=14, bold=True, color=ROUGE)
bullet_box(s, 6.8, 2.15, 5.9, 2.0, [
    "Le TCS3200 mesure la fréquence de sortie",
    "proportionnelle à l'intensité lumineuse réfléchie.",
    "",
    "Filtre S2/S3 = sélection couleur (R, G, B)",
    "Polling toutes les 2 secondes",
    "Si jeton valide → POST /api/jeton_scan.php",
    "",
    "Calibrage : 20+ mesures avec le vrai jeton → min/max"
], fontsize=13, color=BLANC)

code3 = (
    "bool jetonValide() {\n"
    "    uint32_t r = lireCouleur(LOW,  LOW);\n"
    "    uint32_t g = lireCouleur(HIGH, HIGH);\n"
    "    uint32_t b = lireCouleur(LOW,  HIGH);\n"
    "    return (r >= JETON_R_MIN && r <= JETON_R_MAX &&\n"
    "            g >= JETON_G_MIN && g <= JETON_G_MAX &&\n"
    "            b >= JETON_B_MIN && b <= JETON_B_MAX);\n"
    "}"
)
box(s, 6.8, 4.2, 5.9, 2.4, code3, fontsize=12, color=ROUGE,
    bg_color=RGBColor(0x1A, 0x04, 0x04))

footer(s, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Jeu Memory
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 4 — Jeu Memory (HTML5 / JavaScript)",
           "8 paires d'émojis · Fisher-Yates · Score enregistré en BDD")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Mécaniques du jeu", fontsize=14, bold=True, color=VIOLET)
bullet_box(s, 0.6, 2.15, 5.4, 2.5, [
    "• 16 cartes = 8 paires d'émojis",
    "• Mélange Fisher-Yates (O(n), uniforme)",
    "• Clic → retourne la carte",
    "• 2 cartes retournées → vérification match",
    "  ✓ Match    = restent affichées",
    "  ✗ Pas match = se referment après 1 s",
    "• Score calculé sur les mouvements & le temps",
    "• Résultat sauvegardé via enregistrer_partie.php",
], fontsize=12, color=BLANC)

box(s, 0.6, 4.75, 5.4, 0.4, "Dotations possibles", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 0.6, 5.2, 5.4, 1.3, [
    "• Bouteille eau  • Stylo  • Clé USB 16 Go  • Collier Baggio",
    "• Gains définis en BDD (table dotations)"
], fontsize=12, color=GRIS_C)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Algorithme Fisher-Yates", fontsize=14, bold=True, color=VIOLET)
code4 = (
    "// Shuffle O(n) — distribution uniforme\n"
    "function shuffle(array) {\n"
    "    for (let i = array.length - 1; i > 0; i--) {\n"
    "        const j = Math.floor(Math.random() * (i+1));\n"
    "        [array[i], array[j]] = [array[j], array[i]];\n"
    "    }\n"
    "    return array;\n"
    "}\n\n"
    "// État du jeu\n"
    "let flipped = [];   // cartes retournées\n"
    "let matched = 0;    // paires trouvées\n"
    "let moves   = 0;    // compteur coups\n\n"
    "function checkMatch() {\n"
    "    if (flipped[0].emoji === flipped[1].emoji) {\n"
    "        matched++;   // paire !\n"
    "    } else {\n"
    "        setTimeout(retournerCartes, 1000);\n"
    "    }\n"
    "    flipped = [];\n"
    "    moves++;\n"
    "}"
)
box(s, 6.8, 2.2, 5.9, 4.4, code4, fontsize=11, color=VIOLET,
    bg_color=RGBColor(0x0D, 0x0A, 0x20))

footer(s, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Système 2ème Chance
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 5 — Système 2ème Chance",
           "Un QR code sur smartphone pour rejouer gratuitement")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Concept UX", fontsize=14, bold=True, color=OR)
bullet_box(s, 0.6, 2.15, 5.4, 4.3, [
    "1.  Client joue sur la borne → PERD",
    "2.  La borne affiche un QR code",
    "3.  Client scanne avec son smartphone",
    "4.  Ouvre l'app mobile Flutter",
    "5.  Rejoue le même jeu GRATUITEMENT",
    "6.  Résultat enregistré dans la MÊME ligne BDD",
    "",
    "Avantages client :",
    "• Deuxième chance sans retourner en caisse",
    "• Expérience gamifiée et immédiate",
    "",
    "Avantages Adictiz :",
    "• Collecte email/téléphone via app mobile",
    "• Génération de leads qualifiés"
], fontsize=12, color=BLANC)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Implémentation BDD — UPDATE vs INSERT", fontsize=14, bold=True, color=OR)
code5 = (
    "-- 1ère partie : INSERT\n"
    "INSERT INTO parties\n"
    "  (id_client, id_jeu, resultat, date_partie)\n"
    "VALUES (:id, :jeu, 'perdu', NOW());\n\n"
    "-- 2ème chance : UPDATE la MÊME ligne\n"
    "UPDATE parties SET\n"
    "  dc_jeu_id      = :jeu2,\n"
    "  dc_gain_libelle= :gain,\n"
    "  dc_gagne       = :gagne,\n"
    "  dc_score       = :score,\n"
    "  dc_date        = NOW()\n"
    "WHERE id = :id_partie;\n\n"
    "-- Avantages :\n"
    "-- • 1 enregistrement = 2 résultats\n"
    "-- • Pas de jointures complexes\n"
    "-- • Stats faciles à calculer\n"
    "-- • Traçabilité parfaite"
)
box(s, 6.8, 2.2, 5.9, 4.4, code5, fontsize=11, color=OR,
    bg_color=RGBColor(0x1A, 0x14, 0x00))

footer(s, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — ESP32-CAM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Mission 6 — ESP32-CAM — Vidéosurveillance",
           "Streaming MJPEG temps réel — modèle AI Thinker OV2640")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Spécifications & accès", fontsize=14, bold=True, color=ROUGE)
bullet_box(s, 0.6, 2.15, 5.4, 2.5, [
    "• Modèle : AI Thinker — capteur OV2640",
    "• Streaming : MJPEG multipart sur port 80",
    "• URL flux : http://[IP-ESP32-CAM]/stream",
    "• Qualité JPEG configurable",
    "• Surveillance optionnelle de la borne",
    "",
    "Défi : l'ESP32-CAM n'a pas de port USB direct",
    "→ Programmation via ESP8266 en mode pont"
], fontsize=12, color=BLANC)

box(s, 0.6, 4.75, 5.4, 0.35, "Procédure de flash", fontsize=13, bold=True, color=BLANC)
flash = (
    "1.  RST (ESP8266) → GND  (neutraliser le pont)\n"
    "2.  TX/RX ESP8266 → U0T/U0R ESP32-CAM\n"
    "3.  GPIO0 ESP32-CAM → GND  (mode flash)\n"
    "4.  Uploader le firmware\n"
    "5.  Débrancher GPIO0 puis RESET"
)
box(s, 0.6, 5.15, 5.4, 1.4, flash, fontsize=11, color=ROUGE,
    bg_color=RGBColor(0x1A, 0x04, 0x04))

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Extrait firmware streaming MJPEG", fontsize=14, bold=True, color=ROUGE)
code6 = (
    "esp_err_t stream_handler(httpd_req_t* req) {\n"
    "    camera_fb_t* fb  = NULL;\n"
    "    esp_err_t    res = ESP_OK;\n\n"
    "    httpd_resp_set_type(req, STREAM_CONTENT_TYPE);\n\n"
    "    while (true) {\n"
    "        fb = esp_camera_fb_get();\n"
    "        if (!fb) { res = ESP_FAIL; break; }\n\n"
    "        // Envoi boundary + headers\n"
    "        httpd_resp_send_chunk(req,\n"
    "            STREAM_BOUNDARY,\n"
    "            strlen(STREAM_BOUNDARY));\n\n"
    "        // Envoi frame JPEG\n"
    "        httpd_resp_send_chunk(req,\n"
    "            (const char*)fb->buf, fb->len);\n\n"
    "        esp_camera_fb_return(fb);\n"
    "        if (res != ESP_OK) break;\n"
    "    }\n"
    "    return res;\n"
    "}"
)
box(s, 6.8, 2.2, 5.9, 4.4, code6, fontsize=11, color=ROUGE,
    bg_color=RGBColor(0x1A, 0x04, 0x04))

footer(s, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Infrastructure Docker
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Infrastructure — Docker Compose",
           "MySQL 8.0 + PHP 8.2 Apache + phpMyAdmin — déploiement reproductible")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Services Docker", fontsize=14, bold=True, color=VERT)
bullet_box(s, 0.6, 2.15, 5.4, 3.5, [
    "db        — MySQL 8.0",
    "          Volume persistant : ./database",
    "          Port : 3306",
    "",
    "web       — PHP 8.2 Apache",
    "          Port : 80 (HTTP direct, Tailscale → HTTPS)",
    "          Volume : ./projet_ange → /var/www/html",
    "",
    "phpmyadmin — Interface admin BDD",
    "          Port : 8081",
    "          URL  : http://localhost:8081",
], fontsize=12, color=BLANC)

box(s, 0.6, 5.75, 5.4, 0.35, "Commandes utiles", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 0.6, 6.1, 5.4, 0.7, [
    "docker compose up -d     # démarrer",
    "docker compose down      # arrêter",
    "docker compose logs -f   # suivre les logs"
], fontsize=11, color=VERT)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Bénéfices pour le projet", fontsize=14, bold=True, color=VERT)
bullet_box(s, 6.8, 2.15, 5.9, 3.8, [
    "Reproductibilité",
    "  Même environnement dev, test, prod",
    "  Plus de 'ça marche sur ma machine'",
    "",
    "Isolation",
    "  PHP et MySQL séparés dans leurs conteneurs",
    "  Volumes persistants → données survivent au redémarrage",
    "",
    "Scalabilité",
    "  Montée en charge facile",
    "  1000+ scans/jour supportés",
    "",
    "HTTPS via Tailscale",
    "  Apache sert HTTP sur :80",
    "  Tailscale gère le tunnel HTTPS en amont"
], fontsize=12, color=BLANC)

footer(s, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Défis & Solutions
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Défis surmontés",
           "Problèmes rencontrés et solutions apportées")

defis = [
    ("Mémoire IRAM ESP8266 à 92%",
     "Suppression des Serial.println(),\nPROGMEM pour les chaînes longues\n→ RAM réduite à 34.8%"),
    ("Calibrage TCS3200 imprécis",
     "20+ mesures répétées avec le jeton rouge\nPlages min/max stockées dans config.h"),
    ("Anti-rebond RFID",
     "Timer 3 secondes entre deux scans\ndu même UID → évite les doubles jeux"),
    ("2ème chance : UX BDD",
     "UPDATE de la même ligne au lieu\nd'un INSERT → pas de doublons"),
    ("Flash ESP32-CAM sans USB",
     "ESP8266 utilisé comme pont FTDI\nGPIO0→GND = mode flash"),
    ("Sécurité SQL injection",
     "100% Prepared Statements PDO\nAucune concaténation de requête"),
]

positions = [(0.3, 1.6), (4.65, 1.6), (8.95, 1.6),
             (0.3, 4.1), (4.65, 4.1), (8.95, 4.1)]

for (defi, sol), (l, t) in zip(defis, positions):
    rect(s, l, t, 3.9, 2.35, GRIS_F)
    rect(s, l, t, 3.9, 0.5, RGBColor(0x3B, 0x1F, 0x6E))
    box(s, l+0.15, t+0.07, 3.6, 0.38, defi, fontsize=12, bold=True, color=OR)
    box(s, l+0.15, t+0.58, 3.6, 1.65, sol, fontsize=11, color=BLANC)

footer(s, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Résultats & Démos
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
header_bar(s, "Résultats & Démos possibles",
           "Ce qui fonctionne aujourd'hui")

rect(s, 0.4, 1.6, 5.8, 5.2, GRIS_F)
box(s, 0.6, 1.7, 5.4, 0.4, "Bilan des réalisations", fontsize=14, bold=True, color=VERT)
bullet_box(s, 0.6, 2.15, 5.4, 4.3, [
    "✅  RFID RC522 fonctionnel",
    "    Détection badge, envoi UID, LEDs OK",
    "",
    "✅  TCS3200 — détection jeton rouge",
    "    Plages calibrées, polling 2 s",
    "",
    "✅  Jeu Memory complet",
    "    Fisher-Yates, score, enregistrement BDD",
    "",
    "✅  Système 2ème chance",
    "    QR code → app mobile → UPDATE BDD",
    "",
    "✅  ESP32-CAM streaming MJPEG",
    "    Flux temps réel sur /stream",
    "",
    "✅  Sécurité BDD complète",
    "    bcrypt, PDO, CSRF, rate limit, logs"
], fontsize=12, color=BLANC)

rect(s, 6.6, 1.6, 6.3, 5.2, GRIS_F)
box(s, 6.8, 1.7, 5.9, 0.4, "Démos pendant l'oral", fontsize=14, bold=True, color=VERT)
bullet_box(s, 6.8, 2.15, 5.9, 3.8, [
    "1.  Scan badge RFID en live",
    "    Présenter la carte → LED verte → log API",
    "",
    "2.  Insérer le jeton rouge",
    "    Capteur TCS3200 → API → session jeu",
    "",
    "3.  Jouer au Memory sur la borne",
    "    Démonstration complète flip + victoire",
    "",
    "4.  phpMyAdmin",
    "    Montrer les lignes de log en BDD",
    "",
    "5.  ESP32-CAM",
    "    Ouvrir http://[IP]/stream dans un navigateur",
], fontsize=12, color=BLANC)

box(s, 6.8, 6.15, 5.9, 0.5,
    "Prochaine étape : calibrage TCS3200 avec le vrai jeton officiel",
    fontsize=12, color=OR, italic=True)

footer(s, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Conclusion
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NOIR)
rect(s, 0, 0, 13.33, 7.5, NOIR)
rect(s, 0, 0, 13.33, 3.5, VIOLET_D)
rect(s, 0, 3.45, 13.33, 0.12, VIOLET)

box(s, 0.6, 0.3, 12, 0.8, "Merci pour votre attention",
    fontsize=40, bold=True, color=BLANC, align=PP_ALIGN.LEFT)
box(s, 0.6, 1.2, 12, 0.55,
    "BTS CIEL-IR 2025-2026 · Projet BorneInteract · Partenaire Adictiz",
    fontsize=18, color=GRIS_C, italic=True)

rect(s, 0.4, 1.95, 3.8, 1.3, RGBColor(0x3B, 0x1F, 0x6E))
box(s, 0.6, 2.05, 3.5, 0.38, "Mon périmètre", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 0.6, 2.45, 3.5, 0.75, [
    "Sécurité BDD · RFID · TCS3200",
    "Memory · 2ème Chance · ESP32-CAM"
], fontsize=11, color=GRIS_C)

rect(s, 4.6, 1.95, 3.8, 1.3, RGBColor(0x3B, 0x1F, 0x6E))
box(s, 4.8, 2.05, 3.5, 0.38, "Points forts", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 4.8, 2.45, 3.5, 0.75, [
    "Sécurité conforme ANSSI",
    "Firmware optimisé (RAM 34.8%)"
], fontsize=11, color=GRIS_C)

rect(s, 8.8, 1.95, 4.0, 1.3, RGBColor(0x3B, 0x1F, 0x6E))
box(s, 9.0, 2.05, 3.7, 0.38, "À finaliser", fontsize=13, bold=True, color=BLANC)
bullet_box(s, 9.0, 2.45, 3.7, 0.75, [
    "Calibrage TCS3200 jeton officiel",
    "Test RFID conditions réelles"
], fontsize=11, color=GRIS_C)

rect(s, 0.4, 3.8, 12.5, 2.9, GRIS_F)
box(s, 0.6, 3.9, 12, 0.4, "Questions ?", fontsize=22, bold=True, color=VIOLET)
bullet_box(s, 0.6, 4.4, 12, 2.1, [
    "Pourquoi UPDATE et pas INSERT pour la 2ème chance ?",
    "    → 1 enregistrement = 2 résultats · pas de jointures complexes · traçabilité claire",
    "",
    "Comment avez-vous réduit la RAM de l'ESP8266 de 92% à 34.8% ?",
    "    → Suppression des Serial.println() + PROGMEM pour les chaînes + optimisation des délais",
    "",
    "Pourquoi PDO plutôt que mysqli ?",
    "    → PDO supporte multi-BDD, meilleure API, prepared statements natifs, exceptions propres"
], fontsize=12, color=BLANC)

footer(s, 13)

# ── Sauvegarde ──────────────────────────────────────────────────
out = "/home/borne/Desktop/projet_borne_ange3/BorneInteract_Etudiant1.pptx"
prs.save(out)
print(f"Fichier généré : {out}")
